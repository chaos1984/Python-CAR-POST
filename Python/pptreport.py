# -*- coding: utf-8 -*-
'''
Created on Tue Oct 10 2017
@author: wangyj04
'''
import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Cm,Pt,Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_PATTERN


			
class StiffPPT():
#Slid master
	def __init__(self,muban_path,case_path):
		self.case_path = case_path
		self.prs = Presentation(muban_path)
		#Layouts
		self.CoverPage = self.prs.slide_layouts[0]
		self.BlankPage = self.prs.slide_layouts[1]
		self.EndPage = self.prs.slide_layouts[2]
		self.black = RGBColor(0x00,0x00,0x00)
		self.white = RGBColor(0xFF,0xFF,0xFF)
		self.pages = 0
	

	def CoverPageCreate(self,Title,Date,Author):
		#CoverPage
		print '#'*10+'%s' %('COVER PAGE')+'#'*10 
		slide = self.prs.slides.add_slide(self.CoverPage)
		shapes = slide.shapes[3]
		shapes.text = Title
		shapes = slide.shapes[1]
		shapes.text = Author
		shapes = slide.shapes[2]
		shapes.text =  Date
		print '#'*20+'\n'
	# def SummaryPageCreate(self,Title = ' ',Analysis = ' ',Software = ' ',Author = ' ',Date = ' ',AnalysisContent = ' ',AnalysisConclusion = ' '):
		# slide = self.prs.slides.add_slide(self.SummaryPage)
		# shapes = slide.shapes[0]
		# shapes.text = Title
		# shapes = slide.shapes[1]
		# shapes.text = ' '
		# shapes = slide.shapes[2]
		# shapes.text = Analysis
		# shapes = slide.shapes[3]
		# shapes.text = Software
		# shapes = slide.shapes[4]
		# shapes.text = Author
		# shapes = slide.shapes[5]
		# shapes.text = Date
		# shapes = slide.shapes[6]
		# shapes.text = ' '
		# shapes = slide.shapes[7]
		# shapes.text = ' '
		# shapes = slide.shapes[8]
		# shapes.text = ' '
		# shapes = slide.shapes[9]
		# shapes.text = ' '
		# shapes = slide.shapes[10]
		# shapes.text = ' '
		# shapes = slide.shapes[11]
		# shapes.text = ' '
		# shapes = slide.shapes[12]
		# shapes.text = AnalysisContent
		# shapes = slide.shapes[13]
		# shapes.text = AnalysisConclusion
#Content


#Description1
	def BlankPageCreate(self,haeding,subheading,Paragraphs = [],Pictures = [],Movies=[],Tables=[]):
		self.pages += 1
		print '#'*10+'PAGE %d' %(self.pages)+'#'*10 
		slide = self.prs.slides.add_slide(self.BlankPage)
		shapes = slide.shapes[0]
		shapes.text = haeding
		shapes = slide.shapes[1]
		shapes.text = subheading
		tf = shapes.text_frame

		try:
			print 'No. Paragraphs:',len(Paragraphs)
			for i in range(len(Paragraphs)):
				p = tf.add_paragraph()
				p.text = Paragraphs[i]
				p.level = 1
		except:
			print (u'Warnning: paragraphs!')
		

		self.addPictures(slide,Pictures)
		self.addMovies(slide,Movies)
		self.addTables(slide,Tables)
		print '#'*20+'\n'
		
	def EndPageCreate(self):
		self.pages += 1
		print '#'*10+'PAGE %d' %(self.pages)+'#'*10 
		slide = self.prs.slides.add_slide(self.EndPage)
		print '#'*20+'\n'


	def addTables(self,slide,Tables):
		try:
			print 'No. Tables:',len(Tables)
			for i in range(len(Tables)):
				rows = len(Tables[i][0]);cols = len(Tables[i][0][0])
				MatTable = slide.shapes.add_table(rows,cols,Tables[i][1],Tables[i][2],Tables[i][4],Tables[i][3]).table
				# MatTable.vert_banding =True
				MatTable.first_row = False
				for j in range(rows):
					
					for k in range(cols):
						MatTable.cell(j,k).fill.background()
						# MatTable.fore_color.rgb = RGBColor(0x01, 0x23, 0x45)
						MatTable.cell(j,k).text = Tables[i][0][j][k]
						elem_cell = MatTable.cell(j,k)
						elem_cell.fill.solid()
						if j == 0:
							elem_cell.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2
						else:
							elem_cell.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_3
					
						
						# fill = MatTable.cell(j,k).fill
						# fill.patterned()
						# fill.pattern= MSO_PATTERN.PERCENT_90
						# elem_cell.fill.fore_color.rgb = RGBColor(0x01, 0x23, 0x45)
		except:
			print (u'Warnning: Tables!')
			
	
	def addPictures(self,slide,Pictures):
		try:
			print 'No. Pictures:',len(Pictures)
			for i in range(len(Pictures)):
				slide.shapes.add_picture(self.case_path + Pictures[i][0],Pictures[i][1],Pictures[i][2],Pictures[i][4],Pictures[i][3])
		except:
			print (u'Warnning: figures!')

	def addMovies(self,slide,Movies):
		try:
			print 'No. Movies:',len(Movies)
			for i in range(len(Movies)):
				slide.shapes.add_movie(self.case_path + Movies[i][0],Movies[i][1],Movies[i][2],Movies[i][4],Movies[i][3])
		except:
			print (u'Warnning: movies!')	

	def PPTCreate(self,AnalysisType):
		print self.case_path + AnalysisType+'.pptx'
		self.prs.save(self.case_path + AnalysisType+'.pptx')



			
if __name__ == '__main__':
	AnalysisType ='Anchor'
	Title = u'ESR-036622_SX11_AND_SX12_rear_single_anchor_plate_strength'
	Date = '2018-03-09'
	Author = u'Allen Zhu'
	ReportContentsList = ['1.FE Model\n','2.Material Specification\n','3.Load & Boundary Condition\n','4.Results\n','5.Conclusion\n']
	AnalysisContent = 'FE model'
	AnalysisConclusion = 'Current maximum displacement occurs on the anchor plate is 16.7mm, and no breakage occurs on the anchor plate during a tensile force of 22KN.In general, current design meets the analysis requirement.'
	ReportContents = '\n'.join(ReportContentsList)
	case_path = r"C:\Users\yujin.wang\Desktop\New_folder/"
	data_path = r'\n'
	muban_path = r'C:\Users\yujin.wang\Desktop\New_folder\ALV_General Presentation 2017.pptx'


	A = StiffPPT(muban_path,case_path)

	#page1
	A.CoverPageCreate(Title,Date,Author)

	#page2
	A.BlankPageCreate('Contents',ReportContents)

	#Page3
	Pictures = [[r'image\test.png',Cm(1),Cm(10),Cm(8),Cm(15)]]
	TableValue =[[' ','EX\n(MPa)','PRXY','DENS\n(Ton/mm^3)'],['1','2.1E11','0.3','7850']]
	Tables = [[TableValue ,Cm(16),Cm(3),Cm(4),Cm(16)]]
	A.BlankPageCreate('FE Model','',['12','32'],Pictures=Pictures,Tables=Tables)

	#page4
	Pictures = [[r'image\test.png',Cm(16),Cm(4),Cm(8),Cm(16)],[r'image\test.png',Cm(11),Cm(4),Cm(8),Cm(16)]]
	A.BlankPageCreate("Material Specification","Material Type: S500MC\nDensity 7.85e-09 t/mm3 \n Young's modulus 210000 Mpa\n New failure model has been used for this material element will be removed if failure occur",Pictures=Pictures)

	#page5
	Pictures = [[r'image\test.png',Cm(8),Cm(5),Cm(12),Cm(18)]]
	A.BlankPageCreate('Material Specification',"S500MC\nDensity 7.85e-09 t/mm3 \n Young's modulus 210000 Mpa\n New failure model has been used for this material element will be removed if failure occur",Pictures=Pictures)

	#page6
	Pictures = [[r'image\test.png',Cm(15),Cm(5),Cm(6),Cm(13)],[r'image\test.png',Cm(15),Cm(12),Cm(6),Cm(13)]]
	Movies = [[r'image\111.avi',Cm(3),Cm(8),Cm(10),Cm(9)]]
	A.BlankPageCreate('Simulation Results (load 1)','',Pictures=Pictures,Movies=Movies)


	#page7
	load_list =[ ]

	Pictures = [[r'image\test.png',Cm(15),Cm(5),Cm(6),Cm(13)],[r'image\test.png',Cm(15),Cm(12),Cm(6),Cm(13)]]
	Movies = [[r'image\111.avi',Cm(3),Cm(8),Cm(10),Cm(9)]]
	A.BlankPageCreate('Simulation Results (load 2)','',Pictures=Pictures,Movies=Movies)

	#page8
	A.BlankPageCreate('Conclusion','Current maximum displacement occurs on the anchor plate is 16.7mm, and no breakage occurs on the anchor plate during a tensile force of 22KN.\nIn general, current design meets the analysis requirement.')

	A.EndPageCreate()
	A.PPTCreate(AnalysisType)
